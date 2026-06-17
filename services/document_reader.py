"""Evidence-oriented document extraction for summaries and grounded answers."""
from __future__ import annotations

from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
import re

from app.logger import logger


TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".rs", ".go",
    ".java", ".cpp", ".c", ".h", ".cs", ".rb", ".php", ".html", ".css",
    ".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".log",
    ".env", ".sh", ".bat", ".csv", ".sql",
}


@dataclass
class ExtractedContent:
    text: str
    page_count: int = 0
    sampled_pages: list[int] = field(default_factory=list)
    chars_extracted: int = 0
    warnings: list[str] = field(default_factory=list)

    @property
    def is_weak_text(self) -> bool:
        return any("weak text layer" in item.lower() for item in self.warnings)

    def as_prompt_block(self, path: str) -> str:
        target = Path(path)
        lines = [
            f"File: {target.name}",
            f"Type: {target.suffix.lower() or 'file'}",
        ]
        if self.page_count:
            lines.append(f"PDF pages: {self.page_count}")
        if self.sampled_pages:
            lines.append("Sampled pages: " + ", ".join(str(page) for page in self.sampled_pages))
        if self.warnings:
            lines.append("Extraction notes: " + " ".join(self.warnings))
        lines.append("")
        lines.append("Extracted evidence:")
        lines.append(self.text.strip())
        return "\n".join(lines).strip()


def ocr_available() -> bool:
    """Return True when an offline OCR engine is available to Python."""
    try:
        import pytesseract
        from PIL import Image  # noqa: F401
    except Exception:
        return False
    try:
        pytesseract.get_tesseract_version()
    except Exception:
        return False
    return True


def read_file_content(path: str, max_chars: int = 4000) -> str:
    """Compatibility helper that returns only extracted text."""
    return extract_document_sample(path, max_chars=max_chars).text[:max_chars].strip()


def extract_document_sample(path: str, max_chars: int = 12000) -> ExtractedContent:
    """Extract a summary-friendly sample from a file.

    PDF summaries need more than "first N characters". This samples first,
    middle, and last pages, removes repeated boilerplate, and reports when a
    PDF probably needs OCR.
    """
    target = Path(path)
    ext = target.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf_sample(target, max_chars=max_chars)
        if ext in TEXT_EXTENSIONS:
            return ExtractedContent(text=_read_text(target, max_chars), chars_extracted=max_chars)
        if ext in {".docx", ".doc"}:
            return ExtractedContent(text=_read_docx(target, max_chars), chars_extracted=max_chars)
        if ext == ".pptx":
            return ExtractedContent(text=_read_pptx(target, max_chars), chars_extracted=max_chars)
        if ext in {".xlsx", ".xls"}:
            return ExtractedContent(text=_read_xlsx(target, max_chars), chars_extracted=max_chars)
    except Exception as exc:
        logger.warning(f"DocumentReader: cannot extract {path}: {exc}")
        return ExtractedContent(text="", warnings=[f"Extraction failed: {exc}"])
    return ExtractedContent(text="", warnings=["Unsupported or binary file type."])


def _read_text(path: Path, max_chars: int) -> str:
    with path.open("r", encoding="utf-8", errors="ignore") as handle:
        return handle.read(max_chars)


def _sample_page_indexes(page_count: int) -> list[int]:
    if page_count <= 0:
        return []
    if page_count <= 12:
        return list(range(page_count))

    picks = {0, 1, 2, 3, 4, page_count - 3, page_count - 2, page_count - 1}
    for ratio in (0.25, 0.50, 0.75):
        picks.add(min(page_count - 1, max(0, int(page_count * ratio))))
    return sorted(picks)


def _normalize_line(line: str) -> str:
    line = re.sub(r"\s+", " ", line or "").strip()
    line = line.replace("\u200b", "")
    return line


def _is_noise_line(line: str) -> bool:
    if not line:
        return True
    if re.fullmatch(r"\d{1,5}", line):
        return True
    if len(line) <= 2:
        return True
    return False


def _extract_pdf_sample(path: Path, max_chars: int) -> ExtractedContent:
    try:
        import fitz
    except Exception as exc:
        return ExtractedContent(text="", warnings=[f"PyMuPDF unavailable: {exc}"])

    doc = fitz.open(str(path))
    page_count = doc.page_count
    indexes = _sample_page_indexes(page_count)
    page_lines: list[tuple[int, list[str]]] = []
    line_counts: Counter[str] = Counter()
    chars_by_page: list[int] = []

    try:
        for index in indexes:
            text = doc[index].get_text() or ""
            chars_by_page.append(len(text.strip()))
            lines = [
                _normalize_line(line)
                for line in text.splitlines()
                if not _is_noise_line(_normalize_line(line))
            ]
            page_lines.append((index + 1, lines))
            line_counts.update({line.lower(): 1 for line in lines})
    finally:
        doc.close()

    repeated = {line for line, count in line_counts.items() if count >= 3}
    chunks: list[str] = []
    for page_number, lines in page_lines:
        meaningful = [line for line in lines if line.lower() not in repeated]
        if meaningful:
            chunks.append(f"[page {page_number}] " + " ".join(meaningful[:12]))
        elif lines:
            chunks.append(f"[page {page_number}] Only repeated header/footer text was extractable.")

    warnings: list[str] = []
    average_chars = sum(chars_by_page) / max(1, len(chars_by_page))
    if page_count and average_chars < 250:
        warnings.append(
            "Weak text layer detected; this PDF is probably scanned or image-heavy. "
            "OCR is required for a high-confidence summary."
        )
    if repeated:
        warnings.append("Repeated headers/footers were removed from the sampled evidence.")

    evidence = "\n".join(chunks).strip()
    if page_count and average_chars < 250 and ocr_available():
        ocr_text = _ocr_pdf_pages(path, indexes[:4], max_chars=max_chars)
        if ocr_text.strip():
            evidence = ocr_text.strip()
            warnings = [
                warning for warning in warnings
                if "weak text layer" not in warning.lower()
            ]
            warnings.append("OCR was used because the embedded PDF text layer was weak.")
    elif page_count and average_chars < 250:
        warnings.append(
            "Offline OCR engine is not installed, so Neuron cannot read the page images yet."
        )

    if not evidence:
        evidence = "(No meaningful text was extractable from sampled pages.)"
    return ExtractedContent(
        text=evidence[:max_chars],
        page_count=page_count,
        sampled_pages=[index + 1 for index in indexes],
        chars_extracted=len(evidence),
        warnings=warnings,
    )


def _ocr_pdf_pages(path: Path, page_indexes: list[int], max_chars: int) -> str:
    try:
        import fitz
        import pytesseract
        from PIL import Image
    except Exception:
        return ""

    doc = fitz.open(str(path))
    chunks: list[str] = []
    try:
        for index in page_indexes:
            page = doc[index]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(image) or ""
            cleaned = " ".join(text.split())
            if cleaned:
                chunks.append(f"[ocr page {index + 1}] {cleaned}")
            if sum(len(chunk) for chunk in chunks) >= max_chars:
                break
    except Exception as exc:
        logger.warning(f"DocumentReader: PDF OCR failed for {path}: {exc}")
    finally:
        doc.close()
    return "\n".join(chunks)[:max_chars]


def _read_docx(path: Path, max_chars: int) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)[:max_chars]


def _read_pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    text: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)[:max_chars]


def _read_xlsx(path: Path, max_chars: int) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        lines: list[str] = []
        for ws in wb.worksheets[:3]:
            for row in ws.iter_rows(max_row=80, values_only=True):
                values = [str(cell) for cell in row if cell is not None]
                if values:
                    lines.append(" | ".join(values))
                if sum(len(line) for line in lines) >= max_chars:
                    break
        return "\n".join(lines)[:max_chars]
    finally:
        wb.close()
