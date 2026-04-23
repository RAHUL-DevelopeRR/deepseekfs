"""
Neuron — MemoryOS Tool Registry
================================
14 tools for the agentic ReAct loop.
Each tool has: name, description, parameter schema, permission level, execute().

Permission levels:
- SAFE: Read-only operations (auto-execute)
- MODERATE: Write operations (show confirmation dialog)
- DANGEROUS: Destructive operations (block + require explicit override)
"""
from __future__ import annotations

import enum
import json
import os
import re
import glob as glob_module
import hashlib
import shutil
import subprocess
import tempfile
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from app.logger import logger


class PermissionLevel(enum.Enum):
    SAFE = "safe"           # Auto-execute (read-only)
    MODERATE = "moderate"   # Confirm before executing (writes)
    DANGEROUS = "dangerous" # Block + require override (destructive)


@dataclass
class ToolResult:
    """Result returned by a tool execution."""
    success: bool
    output: str
    data: Any = None  # Optional structured data


@dataclass
class ToolParam:
    """Parameter definition for a tool."""
    name: str
    type: str  # "string", "integer", "boolean", "path"
    description: str
    required: bool = True
    default: Any = None


class BaseTool(ABC):
    """Abstract base class for all MemoryOS tools."""
    
    @property
    @abstractmethod
    def name(self) -> str: ...
    
    @property
    @abstractmethod
    def description(self) -> str: ...
    
    @property
    @abstractmethod
    def permission(self) -> PermissionLevel: ...
    
    @property
    @abstractmethod
    def parameters(self) -> List[ToolParam]: ...
    
    @abstractmethod
    def execute(self, **kwargs) -> ToolResult: ...
    
    def to_description_str(self) -> str:
        """Generate a human-readable tool description for the LLM system prompt."""
        params = ", ".join(
            f"{p.name}: {p.type}" + (" (optional)" if not p.required else "")
            for p in self.parameters
        )
        return f"- **{self.name}**({params}): {self.description}"


# ═══════════════════════════════════════════════════════════════
# FILE TOOLS
# ═══════════════════════════════════════════════════════════════

class FileReadTool(BaseTool):
    name = "file_read"
    description = "Read the contents of a file. Returns text content for text files, or OCR for PDFs."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("path", "path", "Absolute path to the file to read"),
        ToolParam("max_chars", "integer", "Maximum characters to read", required=False, default=4000),
    ]
    
    def execute(self, path: str, max_chars: int = 4000, **kwargs) -> ToolResult:
        try:
            if not os.path.isfile(path):
                return ToolResult(False, f"File not found: {path}")
            
            from services.ollama_service import _read_file_content
            content = _read_file_content(path, max_chars=max_chars)
            if content:
                return ToolResult(True, content, {"path": path, "chars": len(content)})
            return ToolResult(False, f"Could not read content from {path}")
        except Exception as e:
            return ToolResult(False, f"Error reading file: {e}")


class FileWriteTool(BaseTool):
    name = "file_write"
    description = "Create a new file with the given content. Supports any text-based format (.py, .html, .json, etc.). Can also create .docx, .xlsx, .pptx via libraries."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("path", "path", "Absolute path for the new file"),
        ToolParam("content", "string", "Content to write to the file"),
        ToolParam("overwrite", "boolean", "If True, overwrite existing file", required=False, default=False),
    ]
    
    def execute(self, path: str, content: str, overwrite: bool = False, **kwargs) -> ToolResult:
        try:
            file_path = Path(path)
            if file_path.exists() and not overwrite:
                return ToolResult(False, f"File already exists: {path}. Set overwrite=True to replace.")
            
            # Create parent directories
            file_path.parent.mkdir(parents=True, exist_ok=True)
            
            ext = file_path.suffix.lower()
            
            # Handle special document formats
            if ext == ".docx":
                return self._write_docx(file_path, content)
            elif ext == ".xlsx":
                return self._write_xlsx(file_path, content)
            elif ext == ".pptx":
                return self._write_pptx(file_path, content)
            else:
                # Standard text file
                with open(path, "w", encoding="utf-8") as f:
                    f.write(content)
                return ToolResult(True, f"Created {path} ({len(content)} chars)", {"path": path})
                
        except Exception as e:
            return ToolResult(False, f"Error writing file: {e}")
    
    def _write_docx(self, path: Path, content: str) -> ToolResult:
        try:
            from docx import Document
            doc = Document()
            for line in content.split("\n"):
                doc.add_paragraph(line)
            doc.save(str(path))
            return ToolResult(True, f"Created Word document: {path}")
        except ImportError:
            return ToolResult(False, "python-docx not installed")
    
    def _write_xlsx(self, path: Path, content: str) -> ToolResult:
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            for i, line in enumerate(content.split("\n"), 1):
                cells = line.split("|") if "|" in line else line.split(",")
                for j, cell in enumerate(cells, 1):
                    ws.cell(row=i, column=j, value=cell.strip())
            wb.save(str(path))
            return ToolResult(True, f"Created Excel file: {path}")
        except ImportError:
            return ToolResult(False, "openpyxl not installed")
    
    def _write_pptx(self, path: Path, content: str) -> ToolResult:
        try:
            from pptx import Presentation
            prs = Presentation()
            for slide_text in content.split("\n---\n"):
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
                lines = slide_text.strip().split("\n")
                if lines:
                    slide.shapes.title.text = lines[0]
                if len(lines) > 1:
                    slide.placeholders[1].text = "\n".join(lines[1:])
            prs.save(str(path))
            return ToolResult(True, f"Created PowerPoint: {path}")
        except ImportError:
            return ToolResult(False, "python-pptx not installed")


class FileEditTool(BaseTool):
    name = "file_edit"
    description = "Edit an existing file by finding and replacing text. Use for modifying code, config files, etc."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("path", "path", "Absolute path to the file to edit"),
        ToolParam("find", "string", "Text to find (exact match)"),
        ToolParam("replace", "string", "Text to replace it with"),
    ]
    
    def execute(self, path: str, find: str, replace: str, **kwargs) -> ToolResult:
        try:
            if not os.path.isfile(path):
                return ToolResult(False, f"File not found: {path}")
            
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            
            if find not in content:
                return ToolResult(False, f"Text not found in {Path(path).name}: '{find[:50]}...'")
            
            new_content = content.replace(find, replace, 1)  # Replace first occurrence
            
            with open(path, "w", encoding="utf-8") as f:
                f.write(new_content)
            
            return ToolResult(True, f"Edited {Path(path).name}: replaced '{find[:40]}' with '{replace[:40]}'")
        except Exception as e:
            return ToolResult(False, f"Error editing file: {e}")


class FileDeleteTool(BaseTool):
    name = "file_delete"
    description = "Delete a file. DANGEROUS — requires explicit user confirmation."
    permission = PermissionLevel.DANGEROUS
    parameters = [
        ToolParam("path", "path", "Absolute path to the file to delete"),
    ]
    
    def execute(self, path: str, **kwargs) -> ToolResult:
        try:
            if not os.path.isfile(path):
                return ToolResult(False, f"File not found: {path}")
            
            os.remove(path)
            return ToolResult(True, f"Deleted: {path}")
        except Exception as e:
            return ToolResult(False, f"Error deleting file: {e}")


# ═══════════════════════════════════════════════════════════════
# FOLDER TOOLS
# ═══════════════════════════════════════════════════════════════

class FolderCreateTool(BaseTool):
    name = "folder_create"
    description = "Create a new folder (and parent directories if needed)."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("path", "path", "Absolute path of the folder to create"),
    ]
    
    def execute(self, path: str, **kwargs) -> ToolResult:
        try:
            folder = Path(path)
            if folder.exists():
                return ToolResult(True, f"Folder already exists: {path}")
            folder.mkdir(parents=True, exist_ok=True)
            return ToolResult(True, f"Created folder: {path}")
        except Exception as e:
            return ToolResult(False, f"Error creating folder: {e}")


class FolderListTool(BaseTool):
    name = "folder_list"
    description = "List contents of a directory (files and subdirectories). Returns a tree-like structure."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("path", "path", "Absolute path of the directory to list"),
        ToolParam("max_depth", "integer", "Maximum depth to traverse", required=False, default=2),
        ToolParam("max_items", "integer", "Maximum items to return", required=False, default=100),
    ]
    
    def execute(self, path: str, max_depth: int = 2, max_items: int = 100, **kwargs) -> ToolResult:
        try:
            folder = Path(path)
            if not folder.is_dir():
                return ToolResult(False, f"Not a directory: {path}")
            
            lines = []
            count = 0
            
            def _walk(dir_path: Path, depth: int, prefix: str):
                nonlocal count
                if depth > max_depth or count >= max_items:
                    return
                try:
                    entries = sorted(dir_path.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower()))
                except PermissionError:
                    lines.append(f"{prefix}[Permission Denied]")
                    return
                
                for entry in entries:
                    if count >= max_items:
                        lines.append(f"{prefix}... (truncated)")
                        return
                    
                    if entry.name.startswith('.'):
                        continue  # Skip hidden files
                    
                    if entry.is_dir():
                        lines.append(f"{prefix}📁 {entry.name}/")
                        count += 1
                        _walk(entry, depth + 1, prefix + "  ")
                    else:
                        size = entry.stat().st_size
                        size_str = _format_size(size)
                        lines.append(f"{prefix}📄 {entry.name} ({size_str})")
                        count += 1
            
            _walk(folder, 0, "")
            output = f"Contents of {path}:\n" + "\n".join(lines) if lines else f"{path} is empty"
            return ToolResult(True, output, {"count": count, "path": path})
        except Exception as e:
            return ToolResult(False, f"Error listing folder: {e}")


class FolderSearchTool(BaseTool):
    name = "folder_search"
    description = "Search for folders by name pattern. Supports glob patterns and fuzzy matching."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("query", "string", "Folder name or pattern to search for"),
        ToolParam("search_path", "path", "Root path to search in", required=False, default=""),
        ToolParam("max_results", "integer", "Max results to return", required=False, default=20),
    ]
    
    def execute(self, query: str, search_path: str = "", max_results: int = 20, **kwargs) -> ToolResult:
        try:
            if not search_path:
                search_path = str(Path.home())
            
            root = Path(search_path)
            if not root.is_dir():
                return ToolResult(False, f"Search path not found: {search_path}")
            
            results = []
            query_lower = query.lower()
            
            for dirpath, dirnames, _ in os.walk(str(root)):
                # Skip hidden and system directories
                dirnames[:] = [d for d in dirnames if not d.startswith('.') and d not in {
                    'node_modules', '__pycache__', '.git', 'venv', '.venv', 'AppData'
                }]
                
                for dirname in dirnames:
                    if query_lower in dirname.lower():
                        full_path = os.path.join(dirpath, dirname)
                        try:
                            file_count = sum(1 for _ in Path(full_path).iterdir())
                        except (PermissionError, OSError):
                            file_count = -1
                        results.append({"path": full_path, "name": dirname, "items": file_count})
                        
                        if len(results) >= max_results:
                            break
                if len(results) >= max_results:
                    break
            
            if not results:
                return ToolResult(True, f"No folders matching '{query}' found in {search_path}")
            
            output_lines = [f"Found {len(results)} folder(s) matching '{query}':"]
            for r in results:
                items = f" ({r['items']} items)" if r['items'] >= 0 else ""
                output_lines.append(f"  📁 {r['path']}{items}")
            
            return ToolResult(True, "\n".join(output_lines), results)
        except Exception as e:
            return ToolResult(False, f"Error searching folders: {e}")


class FolderOrganizeTool(BaseTool):
    name = "folder_organize"
    description = "Organize files in a folder by type, date, or custom categories. Moves files into categorized subfolders."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("path", "path", "Folder to organize"),
        ToolParam("mode", "string", "Organization mode: 'type' (by extension), 'date' (by modified date), 'size' (by size category)"),
        ToolParam("dry_run", "boolean", "If True, only show plan without moving", required=False, default=True),
    ]
    
    # Extension categories for type-based organization
    CATEGORIES = {
        "Documents": {".pdf", ".doc", ".docx", ".txt", ".rtf", ".odt", ".md", ".tex"},
        "Spreadsheets": {".xlsx", ".xls", ".csv", ".ods"},
        "Presentations": {".pptx", ".ppt", ".odp"},
        "Images": {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".webp", ".ico", ".tiff"},
        "Videos": {".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm"},
        "Audio": {".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma"},
        "Archives": {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2"},
        "Code": {".py", ".js", ".ts", ".html", ".css", ".java", ".cpp", ".c", ".rs", ".go", ".rb", ".php"},
        "Executables": {".exe", ".msi", ".bat", ".cmd", ".ps1", ".sh"},
        "Data": {".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".sql", ".db", ".sqlite"},
    }
    
    def execute(self, path: str, mode: str = "type", dry_run: bool = True, **kwargs) -> ToolResult:
        try:
            folder = Path(path)
            if not folder.is_dir():
                return ToolResult(False, f"Not a directory: {path}")
            
            files = [f for f in folder.iterdir() if f.is_file() and not f.name.startswith('.')]
            if not files:
                return ToolResult(True, f"No files to organize in {path}")
            
            # Build the plan
            plan = {}  # category → list of files
            
            if mode == "type":
                for f in files:
                    ext = f.suffix.lower()
                    category = "Other"
                    for cat, exts in self.CATEGORIES.items():
                        if ext in exts:
                            category = cat
                            break
                    plan.setdefault(category, []).append(f)
                    
            elif mode == "date":
                import datetime
                for f in files:
                    mtime = datetime.datetime.fromtimestamp(f.stat().st_mtime)
                    category = mtime.strftime("%Y-%m")
                    plan.setdefault(category, []).append(f)
                    
            elif mode == "size":
                for f in files:
                    size = f.stat().st_size
                    if size < 1024 * 100:  # < 100KB
                        category = "Small (< 100KB)"
                    elif size < 1024 * 1024 * 10:  # < 10MB
                        category = "Medium (100KB - 10MB)"
                    else:
                        category = "Large (> 10MB)"
                    plan.setdefault(category, []).append(f)
            
            # Format plan
            output_lines = [f"Organization plan for {path} (mode: {mode}):"]
            for category, cat_files in sorted(plan.items()):
                output_lines.append(f"\n  📁 {category}/ ({len(cat_files)} files)")
                for f in cat_files[:10]:
                    output_lines.append(f"    → {f.name}")
                if len(cat_files) > 10:
                    output_lines.append(f"    ... and {len(cat_files) - 10} more")
            
            if dry_run:
                output_lines.append(f"\n⚠️ DRY RUN — no files were moved. Set dry_run=False to execute.")
                return ToolResult(True, "\n".join(output_lines), {"plan": {k: [str(f) for f in v] for k, v in plan.items()}, "dry_run": True})
            
            # Execute the plan
            moved = 0
            for category, cat_files in plan.items():
                target_dir = folder / category
                target_dir.mkdir(exist_ok=True)
                for f in cat_files:
                    dest = target_dir / f.name
                    if dest.exists():
                        # Add numeric suffix to avoid overwriting
                        stem = f.stem
                        suffix = f.suffix
                        counter = 1
                        while dest.exists():
                            dest = target_dir / f"{stem}_{counter}{suffix}"
                            counter += 1
                    shutil.move(str(f), str(dest))
                    moved += 1
            
            output_lines.append(f"\n✅ Moved {moved} files into {len(plan)} categories.")
            return ToolResult(True, "\n".join(output_lines), {"moved": moved, "categories": len(plan)})
            
        except Exception as e:
            return ToolResult(False, f"Error organizing folder: {e}")


# ═══════════════════════════════════════════════════════════════
# SEARCH & ANALYSIS TOOLS
# ═══════════════════════════════════════════════════════════════

class SemanticSearchTool(BaseTool):
    name = "semantic_search"
    description = "Search for files using natural language queries. Uses FAISS vector similarity search."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("query", "string", "Natural language search query"),
        ToolParam("max_results", "integer", "Maximum results", required=False, default=10),
    ]
    
    def execute(self, query: str, max_results: int = 10, **kwargs) -> ToolResult:
        try:
            from core.search.semantic_search import SemanticSearch
            searcher = SemanticSearch()
            results = searcher.search(query, top_k=max_results)
            
            if not results:
                return ToolResult(True, f"No results found for: {query}")
            
            output_lines = [f"Search results for '{query}':"]
            for i, r in enumerate(results, 1):
                name = r.get("name", Path(r.get("path", "")).name)
                score = r.get("score", 0)
                output_lines.append(f"  {i}. {name} (score: {score:.2f})")
            
            return ToolResult(True, "\n".join(output_lines), results)
        except Exception as e:
            return ToolResult(False, f"Search error: {e}")


class SummarizeTool(BaseTool):
    name = "summarize"
    description = "Summarize a file or folder contents using AI."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("path", "path", "File or folder path to summarize"),
    ]
    
    def execute(self, path: str, **kwargs) -> ToolResult:
        try:
            p = Path(path)
            if p.is_file():
                from services.llm_engine import get_llm_engine
                summary = get_llm_engine().summarize_file(path)
                return ToolResult(True, f"Summary of {p.name}:\n{summary}")
            elif p.is_dir():
                # Summarize folder structure
                files = list(p.rglob("*"))
                file_count = sum(1 for f in files if f.is_file())
                dir_count = sum(1 for f in files if f.is_dir())
                total_size = sum(f.stat().st_size for f in files if f.is_file())
                
                ext_counts = {}
                for f in files:
                    if f.is_file():
                        ext = f.suffix.lower() or "(no ext)"
                        ext_counts[ext] = ext_counts.get(ext, 0) + 1
                
                top_exts = sorted(ext_counts.items(), key=lambda x: -x[1])[:10]
                
                output = (
                    f"Folder summary: {path}\n"
                    f"  Files: {file_count}\n"
                    f"  Folders: {dir_count}\n"
                    f"  Total size: {_format_size(total_size)}\n"
                    f"  Top file types: {', '.join(f'{ext} ({n})' for ext, n in top_exts)}"
                )
                return ToolResult(True, output)
            else:
                return ToolResult(False, f"Path not found: {path}")
        except Exception as e:
            return ToolResult(False, f"Summarize error: {e}")


class OCRTool(BaseTool):
    name = "ocr"
    description = "Extract text from scanned PDFs and images using OCR. Uses PyMuPDF for digital PDFs."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("path", "path", "Path to PDF or image file"),
    ]
    
    def execute(self, path: str, **kwargs) -> ToolResult:
        try:
            ext = Path(path).suffix.lower()
            
            if ext == ".pdf":
                import fitz
                doc = fitz.open(path)
                text = ""
                for page in doc:
                    text += page.get_text() + "\n"
                doc.close()
                
                if text.strip():
                    return ToolResult(True, f"Extracted text from {Path(path).name}:\n{text[:4000]}")
                return ToolResult(True, "PDF appears to be scanned (no extractable text). OCR not available.")
                
            elif ext in {".jpg", ".jpeg", ".png", ".bmp", ".tiff"}:
                return ToolResult(True, f"Image OCR for {Path(path).name} requires pytesseract (not installed).")
            
            return ToolResult(False, f"Unsupported file type for OCR: {ext}")
        except Exception as e:
            return ToolResult(False, f"OCR error: {e}")


# ═══════════════════════════════════════════════════════════════
# TERMINAL & EXECUTION TOOLS
# ═══════════════════════════════════════════════════════════════

# Dangerous commands that should NEVER be auto-executed
BLOCKED_COMMANDS = {
    "format", "diskpart", "del /s", "rmdir /s", "rd /s",
    "Remove-Item -Recurse -Force", "rm -rf", "shutdown",
    "net user", "reg delete",
}

SAFE_COMMANDS = {
    "dir", "ls", "cd", "pwd", "echo", "type", "cat", "more",
    "where", "which", "whoami", "hostname", "ipconfig",
    "Get-ChildItem", "Get-Item", "Get-Content", "Get-Location",
    "Get-Process", "Get-Service", "Test-Path", "Measure-Object",
    "python --version", "pip --version", "node --version",
    "git status", "git log", "git branch", "git diff",
}


class ShellTool(BaseTool):
    name = "shell"
    description = "Execute a PowerShell or CMD command. Safe commands (dir, ls, cd, etc.) execute immediately. Write commands require confirmation. Destructive commands are blocked."
    permission = PermissionLevel.MODERATE  # Dynamically upgraded to DANGEROUS
    parameters = [
        ToolParam("command", "string", "Command to execute"),
        ToolParam("cwd", "path", "Working directory", required=False, default=""),
        ToolParam("timeout", "integer", "Timeout in seconds", required=False, default=30),
    ]
    
    def _classify_command(self, cmd: str) -> PermissionLevel:
        """Classify a command's risk level."""
        cmd_lower = cmd.lower().strip()
        
        # Check for blocked/dangerous commands
        for blocked in BLOCKED_COMMANDS:
            if blocked.lower() in cmd_lower:
                return PermissionLevel.DANGEROUS
        
        # Check for safe commands
        for safe in SAFE_COMMANDS:
            if cmd_lower.startswith(safe.lower()):
                return PermissionLevel.SAFE
        
        return PermissionLevel.MODERATE
    
    def execute(self, command: str, cwd: str = "", timeout: int = 30, **kwargs) -> ToolResult:
        try:
            # Classify the command
            risk = self._classify_command(command)
            if risk == PermissionLevel.DANGEROUS:
                return ToolResult(
                    False,
                    f"⛔ BLOCKED: '{command}' is classified as DANGEROUS. "
                    f"This command could cause data loss. Execution denied."
                )
            
            # Set working directory
            work_dir = cwd if cwd and os.path.isdir(cwd) else str(Path.home())
            
            logger.info(f"ShellTool: Executing [{risk.value}]: {command}")
            
            result = subprocess.run(
                ["powershell", "-NoProfile", "-Command", command],
                capture_output=True,
                text=True,
                timeout=timeout,
                cwd=work_dir,
                creationflags=subprocess.CREATE_NO_WINDOW if hasattr(subprocess, 'CREATE_NO_WINDOW') else 0,
            )
            
            output = ""
            if result.stdout:
                output += result.stdout
            if result.stderr:
                output += f"\n[stderr] {result.stderr}" if output else result.stderr
            
            output = output.strip()[:5000]  # Limit output size
            
            if result.returncode == 0:
                return ToolResult(True, output or "(command completed with no output)")
            else:
                return ToolResult(False, f"Command exited with code {result.returncode}:\n{output}")
                
        except subprocess.TimeoutExpired:
            return ToolResult(False, f"Command timed out after {timeout}s: {command}")
        except Exception as e:
            return ToolResult(False, f"Shell error: {e}")


class PythonExecTool(BaseTool):
    name = "python_exec"
    description = "Execute Python code in a subprocess. The code is written to a temp file and run with the system Python interpreter."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("code", "string", "Python code to execute"),
        ToolParam("timeout", "integer", "Timeout in seconds", required=False, default=60),
    ]
    
    def execute(self, code: str, timeout: int = 60, **kwargs) -> ToolResult:
        try:
            # Write code to temp file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
                f.write(code)
                temp_path = f.name
            
            try:
                import sys
                python_exe = sys.executable
                
                result = subprocess.run(
                    [python_exe, temp_path],
                    capture_output=True,
                    text=True,
                    timeout=timeout,
                    cwd=str(Path.home()),
                    creationflags=subprocess.CREATE_NO_WINDOW if hasattr(subprocess, 'CREATE_NO_WINDOW') else 0,
                )
                
                output = ""
                if result.stdout:
                    output += result.stdout
                if result.stderr:
                    output += f"\n[stderr] {result.stderr}" if output else result.stderr
                
                output = output.strip()[:5000]
                
                if result.returncode == 0:
                    return ToolResult(True, output or "(script completed with no output)")
                else:
                    return ToolResult(False, f"Script exited with code {result.returncode}:\n{output}")
                    
            finally:
                os.unlink(temp_path)
                
        except subprocess.TimeoutExpired:
            return ToolResult(False, f"Script timed out after {timeout}s")
        except Exception as e:
            return ToolResult(False, f"Python execution error: {e}")


# ═══════════════════════════════════════════════════════════════
# GLOB TOOL
# ═══════════════════════════════════════════════════════════════

class GlobTool(BaseTool):
    name = "glob"
    description = "Find files matching a glob pattern (e.g., '*.py', '**/*.pdf')."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("pattern", "string", "Glob pattern to match"),
        ToolParam("path", "path", "Root directory to search in", required=False, default=""),
        ToolParam("max_results", "integer", "Max results", required=False, default=50),
    ]
    
    def execute(self, pattern: str, path: str = "", max_results: int = 50, **kwargs) -> ToolResult:
        try:
            root = Path(path) if path else Path.home()
            if not root.is_dir():
                return ToolResult(False, f"Directory not found: {path}")
            
            results = []
            for match in root.glob(pattern):
                if len(results) >= max_results:
                    break
                results.append(str(match))
            
            if not results:
                return ToolResult(True, f"No files matching '{pattern}' in {root}")
            
            output = f"Found {len(results)} file(s) matching '{pattern}':\n"
            output += "\n".join(f"  {r}" for r in results)
            return ToolResult(True, output, results)
        except Exception as e:
            return ToolResult(False, f"Glob error: {e}")


# ═══════════════════════════════════════════════════════════════
# UTILITY
# ═══════════════════════════════════════════════════════════════

def _format_size(size_bytes: int) -> str:
    """Format bytes into human-readable size."""
    if size_bytes < 1024:
        return f"{size_bytes}B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f}KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f}MB"
    return f"{size_bytes / (1024 * 1024 * 1024):.1f}GB"


# ═══════════════════════════════════════════════════════════════
# TOOL REGISTRY
# ═══════════════════════════════════════════════════════════════

# All available tools
ALL_TOOLS: Dict[str, BaseTool] = {}

def _register_tools():
    """Register all tools in the registry."""
    global ALL_TOOLS
    tool_classes = [
        FileReadTool, FileWriteTool, FileEditTool, FileDeleteTool,
        FolderCreateTool, FolderListTool, FolderSearchTool, FolderOrganizeTool,
        SemanticSearchTool, SummarizeTool, OCRTool,
        ShellTool, PythonExecTool,
        GlobTool,
    ]
    for cls in tool_classes:
        tool = cls()
        ALL_TOOLS[tool.name] = tool

_register_tools()


def get_tool(name: str) -> Optional[BaseTool]:
    """Get a tool by name."""
    return ALL_TOOLS.get(name)


def get_all_tools() -> Dict[str, BaseTool]:
    """Get all registered tools."""
    return ALL_TOOLS.copy()


def get_tool_descriptions() -> str:
    """Get formatted descriptions of all tools for LLM system prompt."""
    return "\n".join(tool.to_description_str() for tool in ALL_TOOLS.values())


def execute_tool(name: str, **kwargs) -> ToolResult:
    """Execute a tool by name with the given arguments."""
    tool = get_tool(name)
    if tool is None:
        return ToolResult(False, f"Unknown tool: {name}")
    
    try:
        return tool.execute(**kwargs)
    except Exception as e:
        return ToolResult(False, f"Tool '{name}' execution error: {e}")


# ── Type mapping for OpenAI-format schemas ────────────────────
_TYPE_MAP = {
    "string": "string",
    "path": "string",
    "integer": "integer",
    "boolean": "boolean",
}


def get_tool_schemas() -> List[Dict]:
    """Generate OpenAI-format tool schemas for native function calling.
    
    This is the same format that GPT-4 and Claude use.
    llama-cpp-python's `chatml-function-calling` chat format
    accepts these directly via the `tools=` parameter.
    """
    schemas = []
    for tool in ALL_TOOLS.values():
        properties = {}
        required = []
        for p in tool.parameters:
            properties[p.name] = {
                "type": _TYPE_MAP.get(p.type, "string"),
                "description": p.description,
            }
            if p.required:
                required.append(p.name)

        schemas.append({
            "type": "function",
            "function": {
                "name": tool.name,
                "description": tool.description,
                "parameters": {
                    "type": "object",
                    "properties": properties,
                    "required": required,
                },
            },
        })
    return schemas
