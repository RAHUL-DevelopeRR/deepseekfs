"""File-oriented MemoryOS tools."""
from __future__ import annotations

import os
from pathlib import Path

from .base import BaseTool, PermissionLevel, ToolParam, ToolResult


class FileReadTool(BaseTool):
    name = "file_read"
    description = "Read the contents of a file. Returns text content for text files, or OCR for PDFs."
    permission = PermissionLevel.SAFE
    parameters = [
        ToolParam("path", "path", "Absolute path to the file to read"),
        ToolParam("max_chars", "integer", "Maximum characters to read", required=False, default=4000),
    ]

    def execute(self, path: str, max_chars: int = 4000, **kwargs) -> ToolResult:
        try:
            if not os.path.isfile(path):
                return ToolResult(False, f"File not found: {path}")

            from services.ollama_service import _read_file_content

            content = _read_file_content(path, max_chars=max_chars)
            if content:
                return ToolResult(True, content, {"path": path, "chars": len(content)})
            return ToolResult(False, f"Could not read content from {path}")
        except Exception as exc:
            return ToolResult(False, f"Error reading file: {exc}")


class FileWriteTool(BaseTool):
    name = "file_write"
    description = "Create a new file with the given content. Supports any text-based format (.py, .html, .json, etc.). Can also create .docx, .xlsx, .pptx via libraries."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("path", "path", "Absolute path for the new file"),
        ToolParam("content", "string", "Content to write to the file"),
        ToolParam("overwrite", "boolean", "If True, overwrite existing file", required=False, default=False),
    ]

    def execute(self, path: str, content: str, overwrite: bool = False, **kwargs) -> ToolResult:
        try:
            file_path = Path(path)
            if file_path.exists() and not overwrite:
                return ToolResult(False, f"File already exists: {path}. Set overwrite=True to replace.")

            file_path.parent.mkdir(parents=True, exist_ok=True)
            ext = file_path.suffix.lower()

            if ext == ".docx":
                return self._write_docx(file_path, content)
            if ext == ".xlsx":
                return self._write_xlsx(file_path, content)
            if ext == ".pptx":
                return self._write_pptx(file_path, content)

            with open(path, "w", encoding="utf-8") as handle:
                handle.write(content)
            return ToolResult(True, f"Created {path} ({len(content)} chars)", {"path": path})
        except Exception as exc:
            return ToolResult(False, f"Error writing file: {exc}")

    def _write_docx(self, path: Path, content: str) -> ToolResult:
        try:
            from docx import Document

            doc = Document()
            for line in content.split("\n"):
                doc.add_paragraph(line)
            doc.save(str(path))
            return ToolResult(True, f"Created Word document: {path}")
        except ImportError:
            return ToolResult(False, "python-docx not installed")

    def _write_xlsx(self, path: Path, content: str) -> ToolResult:
        try:
            from openpyxl import Workbook

            workbook = Workbook()
            sheet = workbook.active
            for row_index, line in enumerate(content.split("\n"), 1):
                cells = line.split("|") if "|" in line else line.split(",")
                for col_index, cell in enumerate(cells, 1):
                    sheet.cell(row=row_index, column=col_index, value=cell.strip())
            workbook.save(str(path))
            return ToolResult(True, f"Created Excel file: {path}")
        except ImportError:
            return ToolResult(False, "openpyxl not installed")

    def _write_pptx(self, path: Path, content: str) -> ToolResult:
        try:
            from pptx import Presentation

            presentation = Presentation()
            for slide_text in content.split("\n---\n"):
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                lines = slide_text.strip().split("\n")
                if lines:
                    slide.shapes.title.text = lines[0]
                if len(lines) > 1:
                    slide.placeholders[1].text = "\n".join(lines[1:])
            presentation.save(str(path))
            return ToolResult(True, f"Created PowerPoint: {path}")
        except ImportError:
            return ToolResult(False, "python-pptx not installed")


class FileEditTool(BaseTool):
    name = "file_edit"
    description = "Edit an existing file by finding and replacing text. Use for modifying code, config files, etc."
    permission = PermissionLevel.MODERATE
    parameters = [
        ToolParam("path", "path", "Absolute path to the file to edit"),
        ToolParam("find", "string", "Text to find (exact match)"),
        ToolParam("replace", "string", "Text to replace it with"),
    ]

    def execute(self, path: str, find: str, replace: str, **kwargs) -> ToolResult:
        try:
            if not os.path.isfile(path):
                return ToolResult(False, f"File not found: {path}")

            with open(path, "r", encoding="utf-8", errors="ignore") as handle:
                content = handle.read()

            if find not in content:
                return ToolResult(False, f"Text not found in {Path(path).name}: '{find[:50]}...'")

            new_content = content.replace(find, replace, 1)
            with open(path, "w", encoding="utf-8") as handle:
                handle.write(new_content)

            return ToolResult(
                True,
                f"Edited {Path(path).name}: replaced '{find[:40]}' with '{replace[:40]}'",
            )
        except Exception as exc:
            return ToolResult(False, f"Error editing file: {exc}")


class FileDeleteTool(BaseTool):
    name = "file_delete"
    description = "Delete a file. DANGEROUS - requires explicit user confirmation."
    permission = PermissionLevel.DANGEROUS
    parameters = [
        ToolParam("path", "path", "Absolute path to the file to delete"),
    ]

    def execute(self, path: str, **kwargs) -> ToolResult:
        try:
            if not os.path.isfile(path):
                return ToolResult(False, f"File not found: {path}")

            os.remove(path)
            return ToolResult(True, f"Deleted: {path}")
        except Exception as exc:
            return ToolResult(False, f"Error deleting file: {exc}")
