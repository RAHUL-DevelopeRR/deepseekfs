"""
Neuron — AI Intelligence Service (Compatibility Shim)
======================================================
This module was previously the Ollama REST API integration.
It is now a thin compatibility wrapper around LLMEngine
(llama-cpp-python direct GGUF inference).

ALL existing UI code continues to work unchanged.

Migration:
- Old: OllamaService → HTTP → Ollama Server → llama.cpp → model
- New: OllamaService → LLMEngine → llama-cpp-python → model (in-process)
"""
from __future__ import annotations

import os
import hashlib
import threading
from pathlib import Path
from typing import Optional, Dict

from app.logger import logger


# ── File content extraction (preserved from original) ─────────
def _read_file_content(path: str, max_chars: int = 4000) -> str:
    """Extract text from a file for AI analysis."""
    ext = Path(path).suffix.lower()
    text = ""
    try:
        if ext in {".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
                   ".rs", ".go", ".java", ".cpp", ".c", ".h", ".cs",
                   ".rb", ".php", ".html", ".css", ".json", ".xml",
                   ".yaml", ".yml", ".toml", ".ini", ".cfg", ".log",
                   ".env", ".sh", ".bat", ".csv", ".sql"}:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read(max_chars)
        elif ext == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(path)
                for page in doc:
                    text += page.get_text()
                    if len(text) > max_chars: break
                doc.close()
            except Exception:
                pass
        elif ext in {".docx", ".doc"}:
            try:
                from docx import Document
                doc = Document(path)
                text = "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                pass
        elif ext in {".pptx"}:
            try:
                from pptx import Presentation
                prs = Presentation(path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text + "\n"
            except Exception:
                pass
        elif ext in {".xlsx", ".xls"}:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(path, read_only=True, data_only=True)
                for ws in wb.worksheets[:3]:  # first 3 sheets
                    for row in ws.iter_rows(max_row=50, values_only=True):
                        text += " | ".join(str(c) for c in row if c is not None) + "\n"
                wb.close()
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"OllamaService: cannot read {path}: {e}")

    return text[:max_chars].strip()


class OllamaService:
    """Compatibility wrapper — delegates all work to LLMEngine.
    
    Drop-in replacement for the old Ollama REST API-based service.
    All existing UI code calls these same methods without any changes.
    """

    def __init__(self, model: str = ""):
        # model parameter kept for API compatibility but ignored
        # (LLMEngine uses its own model configuration)
        self._engine = None

    def _get_engine(self):
        """Lazy-load the LLM engine."""
        if self._engine is None:
            from services.llm_engine import get_llm_engine
            self._engine = get_llm_engine()
        return self._engine

    # ── Health check ─────────────────────────────────────────
    def is_available(self) -> bool:
        """Check if the AI engine is ready.
        
        Old behavior: checked Ollama REST API
        New behavior: checks if GGUF model is loaded or available
        """
        engine = self._get_engine()
        if engine.is_loaded:
            return True
        # Try to load if not loaded yet
        from services.model_manager import is_llm_model_available
        return is_llm_model_available()

    def _try_auto_start(self) -> bool:
        """Legacy method — no longer needed (no Ollama server to start).
        Kept for compatibility. Now just loads the model."""
        return self.is_available()

    def reset_availability(self):
        """Force re-check on next call."""
        # No-op in new architecture (model state is always known)
        pass

    # ── Pre-warm (load model into RAM) ───────────────────────
    def pre_warm(self):
        """Load model into RAM in a background thread.
        
        Old behavior: sent tiny prompt to Ollama to load model
        New behavior: calls LLMEngine.load_model() directly
        """
        def _warm():
            try:
                engine = self._get_engine()
                logger.info("Encyl: Loading AI model into RAM...")
                t0 = __import__('time').time()
                engine.load_model()
                elapsed = __import__('time').time() - t0
                if engine.is_loaded:
                    logger.info(f"Encyl: Model loaded in {elapsed:.1f}s")
                else:
                    logger.info(f"Encyl: Model load failed: {engine.load_error}")
            except Exception as e:
                logger.info(f"Encyl: Pre-warm failed: {e}")
        threading.Thread(target=_warm, daemon=True).start()

    # ── Public API (unchanged signatures) ─────────────────────

    def summarize_file(self, path: str) -> str:
        """Generate a 2-3 line summary of a file's content."""
        return self._get_engine().summarize_file(path)

    def ask_about_files(self, question: str, file_contexts: list[dict]) -> str:
        """Answer a natural language question using file context."""
        return self._get_engine().ask_about_files(question, file_contexts)

    def suggest_tags(self, path: str) -> list[str]:
        """Auto-generate tags for a file."""
        return self._get_engine().suggest_tags(path)

    @property
    def cache_size(self) -> int:
        return self._get_engine().cache_size


# ── Singleton ────────────────────────────────────────────────
_instance: Optional[OllamaService] = None

def get_ollama() -> OllamaService:
    global _instance
    if _instance is None:
        _instance = OllamaService()
    return _instance
