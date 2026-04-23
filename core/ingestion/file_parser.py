"""Extract text from various file formats"""
import os
from pathlib import Path
from typing import Optional
import app.config as config
from app.logger import logger

class FileParser:
    """Parse different file types and extract text"""

    @staticmethod
    def parse(file_path: str) -> Optional[str]:
        """Extract text from file based on extension."""
        try:
            ext = Path(file_path).suffix.lower()

            # Documents
            if ext in (".txt", ".md"):
                return FileParser._parse_text(file_path)
            elif ext == ".pdf":
                return FileParser._parse_pdf(file_path)
            elif ext in (".docx", ".doc"):
                return FileParser._parse_docx(file_path)
            elif ext in (".pptx", ".ppt"):
                return FileParser._parse_pptx(file_path)
            elif ext in (".xlsx", ".xls"):
                return FileParser._parse_xlsx(file_path)

            # Data / markup
            elif ext == ".json":
                return FileParser._parse_text(file_path)
            elif ext == ".csv":
                return FileParser._parse_csv(file_path)
            elif ext in (".html", ".htm"):
                return FileParser._parse_html(file_path)
            elif ext == ".xml":
                return FileParser._parse_html(file_path)  # reuse tag stripping

            # Notebooks
            elif ext == ".ipynb":
                return FileParser._parse_notebook(file_path)

            # Logs
            elif ext == ".log":
                return FileParser._parse_log(file_path)

            # Config files
            elif ext in (".env", ".ini", ".toml", ".cfg", ".yaml", ".yml"):
                return FileParser._parse_config(file_path)

            # Code (all languages)
            elif ext in (".py", ".js", ".ts", ".jsx", ".tsx",
                         ".rs", ".go", ".java", ".cpp", ".c", ".h",
                         ".cs", ".rb", ".php", ".swift", ".kt",
                         ".css"):
                return FileParser._parse_text(file_path)

            # Media (metadata from filename only)
            elif ext in (".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm"):
                return FileParser._parse_video_metadata(file_path)

            else:
                return None
        except Exception as e:
            logger.warning(f"Error parsing {file_path}: {e}")
            return None

    # ── Plain text ────────────────────────────────────────────
    @staticmethod
    def _parse_text(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()[:5000]

    # ── PDF ───────────────────────────────────────────────────
    @staticmethod
    def _parse_pdf(file_path: str) -> Optional[str]:
        try:
            import fitz
            doc = fitz.open(file_path)
            text = ""
            for page in doc[:5]:
                text += page.get_text()
            return text[:5000]
        except Exception as e:
            logger.debug(f"PDF parse error: {e}")
            return None

    # ── DOCX ──────────────────────────────────────────────────
    @staticmethod
    def _parse_docx(file_path: str) -> Optional[str]:
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            return text[:5000]
        except Exception as e:
            logger.debug(f"DOCX parse error: {e}")
            return None

    # ── PPTX ──────────────────────────────────────────────────
    @staticmethod
    def _parse_pptx(file_path: str) -> Optional[str]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)[:5000]
        except Exception as e:
            logger.debug(f"PPTX parse error: {e}")
            return None

    # ── XLSX ──────────────────────────────────────────────────
    @staticmethod
    def _parse_xlsx(file_path: str) -> Optional[str]:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            cells = []
            for sheet in wb.sheetnames[:3]:
                ws = wb[sheet]
                count = 0
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            cells.append(str(cell))
                            count += 1
                            if count >= 100:
                                break
                    if count >= 100:
                        break
            wb.close()
            return " ".join(cells)[:5000]
        except Exception as e:
            logger.debug(f"XLSX parse error: {e}")
            return None

    # ── CSV ───────────────────────────────────────────────────
    @staticmethod
    def _parse_csv(file_path: str) -> str:
        import csv
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i >= 50:
                        break
                    text += " ".join(row) + "\n"
            return text[:5000]
        except Exception as e:
            logger.debug(f"CSV parse error: {e}")
            return FileParser._parse_text(file_path)

    # ── HTML / XML ────────────────────────────────────────────
    @staticmethod
    def _parse_html(file_path: str) -> Optional[str]:
        try:
            from html.parser import HTMLParser

            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.texts = []
                    self._skip = False

                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style", "noscript"):
                        self._skip = True

                def handle_endtag(self, tag):
                    if tag in ("script", "style", "noscript"):
                        self._skip = False

                def handle_data(self, data):
                    if not self._skip:
                        stripped = data.strip()
                        if stripped:
                            self.texts.append(stripped)

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()[:50000]

            parser = TextExtractor()
            parser.feed(content)
            return " ".join(parser.texts)[:5000]
        except Exception as e:
            logger.debug(f"HTML parse error: {e}")
            return None

    # ── Jupyter Notebook ──────────────────────────────────────
    @staticmethod
    def _parse_notebook(file_path: str) -> Optional[str]:
        try:
            import json
            with open(file_path, 'r', encoding='utf-8') as f:
                nb = json.load(f)

            texts = []
            for cell in nb.get("cells", []):
                cell_type = cell.get("cell_type", "")
                source = "".join(cell.get("source", []))
                if cell_type in ("markdown", "code") and source.strip():
                    texts.append(source)

            return "\n\n".join(texts)[:8000]
        except Exception as e:
            logger.debug(f"Notebook parse error: {e}")
            return None

    # ── Log files (last 200 lines) ────────────────────────────
    @staticmethod
    def _parse_log(file_path: str) -> Optional[str]:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = f.readlines()
            last_200 = lines[-200:] if len(lines) > 200 else lines
            return "".join(last_200)[:5000]
        except Exception as e:
            logger.debug(f"Log parse error: {e}")
            return None

    # ── Config files (.env, .ini, .toml, .cfg) ────────────────
    @staticmethod
    def _parse_config(file_path: str) -> Optional[str]:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()[:2000]
        except Exception as e:
            logger.debug(f"Config parse error: {e}")
            return None

    # ── Video (filename metadata only) ────────────────────────
    @staticmethod
    def _parse_video_metadata(file_path: str) -> str:
        name = os.path.basename(file_path)
        name_clean = name.replace("_", " ").replace("-", " ").replace(".", " ")
        return f"Video file: {name_clean}"

    # ── File metadata ─────────────────────────────────────────
    @staticmethod
    def get_file_metadata(file_path: str) -> dict:
        stat = os.stat(file_path)
        return {
            "path": file_path,
            "name": os.path.basename(file_path),
            "size": stat.st_size,
            "modified_time": stat.st_mtime,
            "created_time": stat.st_ctime,
            "extension": Path(file_path).suffix.lower(),
        }


def extract_text(file_path: str) -> Optional[str]:
    """Compatibility wrapper for callers/tests that expect extract_text."""
    return FileParser.parse(file_path)
